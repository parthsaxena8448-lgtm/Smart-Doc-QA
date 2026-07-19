import fitz  # PyMuPDF for PDFs
import docx  # For Word documents
from pptx import Presentation  # For PowerPoint files
import os

def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file."""
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    doc.close()
    return text

def extract_text_from_docx(file_path):
    """Extracts text from a Word (.docx) file."""
    doc = docx.Document(file_path)
    text = []
    # Loop through every paragraph in the Word document
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)

def extract_text_from_pptx(file_path):
    """Extracts text from a PowerPoint (.pptx) file."""
    prs = Presentation(file_path)
    text = []
    # Loop through every slide, and then every text box (shape) on that slide
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text(file_path):
    """
    Master function that checks the file type and calls the right extractor.
    """
    # Find out if it's a .pdf, .docx, or .pptx
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        return f"Unsupported file type: {file_extension}"

# --- Testing Section ---
if __name__ == "__main__":
    # You can change this path to test a .docx or .pptx file later!
    test_file_path = "data/test.pdf" 
    
    print(f"Extracting text from: {test_file_path}...\n")
    text_result = extract_text(test_file_path)
    
    print(text_result[:500])
    print("\n\n--- Extraction Complete! ---")